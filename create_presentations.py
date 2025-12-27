#!/usr/bin/env python3
"""
Script to create PowerPoint presentations from Jupyter notebook outputs.
Each assignment gets its own PPTX file with curated images.
"""

import nbformat
import os
import base64
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import hashlib

def calculate_image_hash(img):
    """Calculate a perceptual hash of an image for similarity detection."""
    # Resize to small size for comparison
    small = img.resize((8, 8), Image.Resampling.LANCZOS).convert('L')
    pixels = list(small.getdata())
    avg = sum(pixels) / len(pixels)
    # Create hash based on whether pixel is above/below average
    bits = ''.join('1' if p > avg else '0' for p in pixels)
    return bits

def are_images_similar(img1, img2, threshold=0.90):
    """Check if two images are similar based on hash comparison."""
    hash1 = calculate_image_hash(img1)
    hash2 = calculate_image_hash(img2)
    
    # Calculate similarity (Hamming distance)
    matches = sum(b1 == b2 for b1, b2 in zip(hash1, hash2))
    similarity = matches / len(hash1)
    
    return similarity > threshold

def extract_unique_images_from_notebook(notebook_path, max_similar=2, aggressive=False):
    """
    Extract unique images from a Jupyter notebook with aggressive deduplication.
    
    Args:
        notebook_path: Path to the .ipynb file
        max_similar: Maximum number of similar images to keep per group
        aggressive: If True, applies more aggressive filtering (keeps only 1 per similar group)
        
    Returns:
        List of PIL Image objects
    """
    images = []
    
    try:
        nb = nbformat.read(notebook_path, as_version=4)
        
        # Group images by cell first to understand context
        cell_images = {}
        for cell_idx, cell in enumerate(nb.cells):
            if cell.cell_type == 'code' and hasattr(cell, 'outputs'):
                cell_imgs = []
                for output_idx, output in enumerate(cell.outputs):
                    if hasattr(output, 'data') and 'image/png' in output.data:
                        # Decode base64 image data
                        image_data = output.data['image/png']
                        image_bytes = base64.b64decode(image_data)
                        
                        # Convert to PIL Image
                        img = Image.open(io.BytesIO(image_bytes))
                        cell_imgs.append(img)
                
                if cell_imgs:
                    cell_images[cell_idx] = cell_imgs
        
        # Process each cell's images
        all_images = []
        for cell_idx, imgs in cell_images.items():
            if len(imgs) == 1:
                all_images.append(imgs[0])
            else:
                # Multiple images in same cell - likely variations
                # If aggressive, keep only first; otherwise keep first and last
                if aggressive:
                    all_images.append(imgs[0])
                    print(f"  Cell {cell_idx}: Kept 1 of {len(imgs)} images")
                else:
                    # Check if they're all similar
                    all_similar = all(are_images_similar(imgs[0], imgs[i], threshold=0.88) for i in range(1, len(imgs)))
                    if all_similar and len(imgs) > 3:
                        # Many similar images - keep first and last
                        all_images.append(imgs[0])
                        all_images.append(imgs[-1])
                        print(f"  Cell {cell_idx}: Reduced {len(imgs)} similar images to 2")
                    else:
                        # Not all similar or small group
                        all_images.extend(imgs[:min(len(imgs), max_similar)])
        
        # Additional global deduplication across cells
        threshold = 0.85 if aggressive else 0.88
        similar_groups = []
        for img in all_images:
            added = False
            for group in similar_groups:
                if are_images_similar(img, group[0], threshold=threshold):
                    group.append(img)
                    added = True
                    break
            if not added:
                similar_groups.append([img])
        
        # Keep only one representative from each global group
        for group in similar_groups:
            if len(group) == 1:
                images.append(group[0])
            else:
                images.append(group[0])  # Keep only the first
                if len(group) > 1:
                    print(f"  Globally deduplicated {len(group)} similar images to 1")
        
        print(f"  Final: {len(images)} unique images from {sum(len(imgs) for imgs in cell_images.values())} total")
    
    except Exception as e:
        print(f"  Error reading {notebook_path}: {e}")
    
    return images

def create_presentation(images, output_path, title):
    """
    Create a PowerPoint presentation from images.
    
    Args:
        images: List of PIL Image objects
        output_path: Path to save the PPTX file
        title: Title for the presentation
    """
    if not images:
        print(f"  No images to create presentation at {output_path}")
        return
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title_shape.text = title
    subtitle.text = f"{len(images)} Output Images"
    
    # Add each image on a separate slide
    for idx, img in enumerate(images):
        # Use blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Save image to bytes
        img_bytes = io.BytesIO()
        if img.mode == 'RGBA':
            # Convert RGBA to RGB
            rgb_img = Image.new('RGB', img.size, (255, 255, 255))
            rgb_img.paste(img, mask=img.split()[3])
            rgb_img.save(img_bytes, format='PNG')
        else:
            img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        # Calculate image dimensions to fit slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Leave margins
        margin = Inches(0.5)
        max_width = slide_width - 2 * margin
        max_height = slide_height - 2 * margin
        
        # Calculate scaling
        img_width, img_height = img.size
        width_ratio = max_width / img_width
        height_ratio = max_height / img_height
        scale = min(width_ratio, height_ratio)
        
        final_width = int(img_width * scale)
        final_height = int(img_height * scale)
        
        # Center the image
        left = (slide_width - final_width) / 2
        top = (slide_height - final_height) / 2
        
        # Add image to slide
        slide.shapes.add_picture(img_bytes, left, top, final_width, final_height)
    
    prs.save(output_path)
    print(f"  Created presentation: {output_path} ({len(images)} slides)")

def main():
    """Main function to process all notebooks and create presentations."""
    
    print("=== Creating PowerPoint presentations from notebook outputs ===\n")
    
    # Assignment directories to process
    assignments = {
        'AS1': {
            'notebooks': ['AS1.ipynb'],
            'title': 'Assignment 1: Image Formation & Geometric Transformations'
        },
        'AS2': {
            'notebooks': ['AS2.ipynb'],
            'title': 'Assignment 2: Image Filtering & Edge Detection'
        },
        'AS3': {
            'notebooks': ['AS3_P1.ipynb', 'AS3_P2.ipynb', 'AS3_P3.ipynb'],
            'title': 'Assignment 3: Robust Estimation & CNN Classification'
        },
        'AS4': {
            'notebooks': ['AS4_Q1.ipynb', 'AS4_Q2.ipynb'],
            'title': 'Assignment 4: Semantic Segmentation & Object Detection'
        },
        'AS5': {
            'notebooks': ['AS5.ipynb'],
            'title': 'Assignment 5: Vision Transformers'
        }
    }
    
    presentation_counts = {}
    all_images_combined = []
    
    # Process each assignment
    for assignment, config in assignments.items():
        print(f"Processing {assignment}...")
        assignment_images = []
        
        # Use aggressive filtering for AS4 which has many training graphs
        aggressive = (assignment == 'AS4')
        
        for notebook in config['notebooks']:
            notebook_path = os.path.join(assignment, notebook)
            if os.path.exists(notebook_path):
                print(f"  Reading {notebook_path}...")
                images = extract_unique_images_from_notebook(notebook_path, max_similar=2, aggressive=aggressive)
                assignment_images.extend(images)
            else:
                print(f"  Notebook not found: {notebook_path}")
        
        # Add to combined collection
        all_images_combined.extend(assignment_images)
        
        # Create presentation for this assignment if there are images
        if assignment_images:
            pptx_path = f"{assignment}_outputs.pptx"
            create_presentation(assignment_images, pptx_path, config['title'])
            presentation_counts[assignment] = len(assignment_images)
        else:
            print(f"  No images found in {assignment}")
        
        print()
    
    # Create combined presentation with all images
    print("Creating combined presentation with all images...")
    if all_images_combined:
        create_presentation(all_images_combined, "all_outputs.pptx", 
                          "All Assignments: Computer Vision Outputs")
        print(f"  Combined presentation: all_outputs.pptx ({len(all_images_combined)} slides)")
    else:
        print("No images found across all assignments")
    
    # Print summary
    print("\n=== Summary ===")
    print(f"Total presentations created: {len(presentation_counts) + 1}")
    for assignment, count in presentation_counts.items():
        print(f"  {assignment}_outputs.pptx: {count} slides")
    if all_images_combined:
        print(f"  all_outputs.pptx: {len(all_images_combined)} slides (combined)")
    
    print("\nNote: Similar/repetitive images have been aggressively filtered.")
    print("You can now extract GIFs from these presentations as needed.")

if __name__ == '__main__':
    main()
